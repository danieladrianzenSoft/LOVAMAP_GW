#!/usr/bin/env python3
"""
Generate LOVAMAP Advisory Board presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──────────────────────────────────────────────
DUKE_BLUE    = RGBColor(0x00, 0x36, 0x80)
DUKE_NAVY    = RGBColor(0x01, 0x2A, 0x5E)
ACCENT_TEAL  = RGBColor(0x00, 0x9C, 0xAB)
ACCENT_GREEN = RGBColor(0x3A, 0xA1, 0x7E)
ACCENT_AMBER = RGBColor(0xE8, 0x9C, 0x23)
ACCENT_CORAL = RGBColor(0xE0, 0x5A, 0x4F)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF4, 0xF5, 0xF7)
MID_GRAY     = RGBColor(0x6B, 0x7B, 0x8D)
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
SOFT_WHITE   = RGBColor(0xFA, 0xFA, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ────────────────────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=16, color=DARK_TEXT, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_list(tf, items, size=15, color=DARK_TEXT, indent_level=0, bold_prefix=True, space=Pt(3)):
    """Add items as bullet points. Items can be str or (bold_part, rest) tuples."""
    for item in items:
        p = tf.add_paragraph()
        p.space_before = space
        p.space_after = Pt(1)
        p.level = indent_level

        if isinstance(item, tuple):
            bold_part, rest = item
            run = p.add_run()
            run.text = bold_part
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = True
            run.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = f"  {item}"
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = "Calibri"


def make_pillar_card(slide, left, top, width, height, accent_color, number, title, items):
    """Create a styled pillar card with accent top bar."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, WHITE)
    # Accent bar at top
    add_rect(slide, left + Inches(0.15), top, width - Inches(0.3), Inches(0.06), accent_color)

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.25), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent_color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.word_wrap = False
    set_text(tf, str(number), size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    tb = add_text_box(slide, left + Inches(0.95), top + Inches(0.2), width - Inches(1.2), Inches(0.55))
    set_text(tb.text_frame, title, size=17, color=DUKE_NAVY, bold=True)

    # Bullets
    tb2 = add_text_box(slide, left + Inches(0.3), top + Inches(0.75), width - Inches(0.6), height - Inches(0.9))
    tb2.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tb2.text_frame.paragraphs[0]
            first = False
        else:
            p = tb2.text_frame.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(1)

        if isinstance(item, tuple):
            bold_part, rest = item
            r = p.add_run()
            r.text = bold_part
            r.font.size = Pt(13)
            r.font.color.rgb = DARK_TEXT
            r.font.bold = True
            r.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(13)
            r2.font.color.rgb = MID_GRAY
            r2.font.bold = False
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = f"\u2022  {item}"
            r.font.size = Pt(13)
            r.font.color.rgb = DARK_TEXT
            r.font.name = "Calibri"


def add_section_header(slide, number, title, accent_color):
    """Add a pillar section header bar at the top of a slide."""
    # Accent stripe
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), accent_color)
    # Number badge
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.4), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent_color
    circ.line.fill.background()
    tf = circ.text_frame
    set_text(tf, str(number), size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    tb = add_text_box(slide, Inches(1.4), Inches(0.35), Inches(10), Inches(0.7))
    set_text(tb.text_frame, title, size=28, color=DUKE_NAVY, bold=True)


def add_slide_number(slide, num):
    tb = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb.text_frame, str(num), size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, DUKE_NAVY)

# Subtle geometric accent
add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.08), ACCENT_TEAL)

# Title text
tb = add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(1.2))
set_text(tb.text_frame, "LOVAMAP", size=54, color=WHITE, bold=True)
add_paragraph(tb.text_frame, "Lowering the Barrier for Void Analysis of", size=26, color=RGBColor(0xB0, 0xC4, 0xDE), space_before=Pt(4))
add_paragraph(tb.text_frame, "Microporous Annealed Particle Scaffolds", size=26, color=RGBColor(0xB0, 0xC4, 0xDE), space_before=Pt(0))

# Subtitle
tb2 = add_text_box(slide, Inches(1.0), Inches(4.0), Inches(11), Inches(0.8))
set_text(tb2.text_frame, "Advisory Board Meeting  \u2014  Platform Overview & Discussion", size=20, color=ACCENT_TEAL)

# Attribution
tb3 = add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11), Inches(1.0))
set_text(tb3.text_frame, "Segura Lab  |  Duke University", size=16, color=RGBColor(0x8A, 0x9B, 0xAE))
add_paragraph(tb3.text_frame, "June 2026", size=14, color=RGBColor(0x8A, 0x9B, 0xAE), space_before=Pt(8))


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), DUKE_BLUE)
add_slide_number(slide, 2)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "The Problem", size=32, color=DUKE_NAVY, bold=True)

# Three problem cards
problems = [
    ("Inaccessible Analysis", "ACCENT_CORAL",
     "Running LOVAMAP scaffold analysis currently requires local installation of specialized software, command-line expertise, and significant compute resources \u2014 putting it out of reach for most experimental labs.",
     ACCENT_CORAL),
    ("Fragmented Data", "ACCENT_AMBER",
     "Granular material scaffold data is scattered across individual labs, stored in incompatible formats, and published without standardized metadata \u2014 making cross-study comparison nearly impossible.",
     ACCENT_AMBER),
    ("Inconsistent Reporting", "ACCENT_TEAL",
     "Published studies on MAP scaffolds report fabrication parameters, experimental conditions, and outcomes using varying terminology, units, and levels of detail \u2014 hindering reproducibility.",
     ACCENT_TEAL),
]

card_w = Inches(3.7)
card_h = Inches(4.5)
gap = Inches(0.45)
start_x = Inches(0.65)

for i, (title, _, desc, accent) in enumerate(problems):
    left = start_x + i * (card_w + gap)
    top = Inches(1.6)
    card = add_rounded_rect(slide, left, top, card_w, card_h, WHITE)
    add_rect(slide, left + Inches(0.1), top, card_w - Inches(0.2), Inches(0.06), accent)

    ttb = add_text_box(slide, left + Inches(0.35), top + Inches(0.35), card_w - Inches(0.7), Inches(0.5))
    set_text(ttb.text_frame, title, size=19, color=accent, bold=True)

    dtb = add_text_box(slide, left + Inches(0.35), top + Inches(1.0), card_w - Inches(0.7), card_h - Inches(1.3))
    dtb.text_frame.word_wrap = True
    set_text(dtb.text_frame, desc, size=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — OUR APPROACH: THREE PILLARS
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), DUKE_BLUE)
add_slide_number(slide, 3)

tb = add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Our Approach: Three Pillars", size=32, color=DUKE_NAVY, bold=True)

sub = add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
set_text(sub.text_frame, "LOVAMAP Gateway addresses these challenges through an integrated web platform built on three pillars.", size=15, color=MID_GRAY)

card_w = Inches(3.7)
card_h = Inches(4.3)
start_x = Inches(0.65)
top = Inches(1.9)

make_pillar_card(slide, start_x, top, card_w, card_h, ACCENT_TEAL, 1,
    "Scaffold Analysis",
    [
        "\u2022  Submit jobs to lovamap_core via browser",
        "\u2022  Particle segmentation from microscopy",
        "\u2022  Pore-level structural descriptors",
        "\u2022  3D mesh generation & visualization",
        "\u2022  Job chaining (segment \u2192 analyze \u2192 mesh)",
        "\u2022  AI-powered scaffold search",
    ])

make_pillar_card(slide, start_x + card_w + Inches(0.45), top, card_w, card_h, ACCENT_GREEN, 2,
    "Data Storage",
    [
        "\u2022  PostgreSQL: geometry & structure",
        "\u2022  RDF/Fuseki: literature & experiments",
        "\u2022  Per-pore descriptor arrays",
        "\u2022  3D mesh files (.glb) for each domain",
        "\u2022  Paper \u2192 Material \u2192 Experiment \u2192 Outcome",
        "\u2022  Bridge: RDF links to relational scaffolds",
    ])

make_pillar_card(slide, start_x + 2 * (card_w + Inches(0.45)), top, card_w, card_h, ACCENT_AMBER, 3,
    "Data Standardization",
    [
        "\u2022  Controlled descriptor vocabulary",
        "\u2022  Structured fabrication metadata",
        "\u2022  RDF ontology as reporting schema",
        "\u2022  Publication dataset versioning",
        "\u2022  Immutable, hash-verified snapshots",
        "\u2022  Auto-tagging & classification",
    ])


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), DUKE_BLUE)
add_slide_number(slide, 4)

tb = add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.8))
set_text(tb.text_frame, "System Architecture", size=32, color=DUKE_NAVY, bold=True)

# ── Top row: Frontend ──
box_h = Inches(0.9)
fe_w = Inches(8)
fe_left = Inches(2.65)
fe_top = Inches(1.5)
fe = add_rounded_rect(slide, fe_left, fe_top, fe_w, box_h, RGBColor(0xE8, 0xEA, 0xF0))
ttb = add_text_box(slide, fe_left, fe_top + Inches(0.05), fe_w, box_h)
ttb.text_frame.word_wrap = True
set_text(ttb.text_frame, "LOVAMAP Frontend", size=16, color=DUKE_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(ttb.text_frame, "React + TypeScript  |  3D Visualization  |  Knowledge Graph", size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down
arr_tb = add_text_box(slide, Inches(6.2), fe_top + box_h - Inches(0.05), Inches(1), Inches(0.45))
set_text(arr_tb.text_frame, "\u25BC", size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ── Middle row: Gateway ──
gw_top = Inches(2.8)
gw_h = Inches(2.6)
gw = add_rounded_rect(slide, fe_left, gw_top, fe_w, gw_h, RGBColor(0xE0, 0xF0, 0xF8))
add_rect(slide, fe_left + Inches(0.1), gw_top, fe_w - Inches(0.2), Inches(0.05), DUKE_BLUE)

ttb2 = add_text_box(slide, fe_left, gw_top + Inches(0.15), fe_w, Inches(0.4))
set_text(ttb2.text_frame, "LOVAMAP Gateway  (.NET 8)", size=16, color=DUKE_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Service boxes inside gateway
svc_labels = [
    ("Jobs &\nCore Integration", ACCENT_TEAL),
    ("Scaffold\nGroups", ACCENT_GREEN),
    ("Descriptors\n& Domains", ACCENT_GREEN),
    ("RDF\nService", ACCENT_AMBER),
    ("Publications\n& Datasets", ACCENT_AMBER),
    ("Auth &\nUsers", MID_GRAY),
]
svc_w = Inches(1.15)
svc_h = Inches(1.05)
svc_gap = Inches(0.12)
svc_start_x = fe_left + Inches(0.25)
svc_top = gw_top + Inches(0.65)
# Two rows of 3
for row in range(2):
    for col in range(3):
        idx = row * 3 + col
        label, color = svc_labels[idx]
        sx = svc_start_x + col * (svc_w + Inches(0.08)) + row * Inches(3.85)
        # For second row, reset
        if row == 1:
            sx = svc_start_x + col * (svc_w + Inches(0.08)) + Inches(3.85)
# Actually let me just lay them out in a single row of 6
svc_w = Inches(1.15)
total_svc_w = 6 * svc_w + 5 * svc_gap
svc_start_x = fe_left + (fe_w - total_svc_w) / 2
svc_top = gw_top + Inches(0.7)
for i, (label, color) in enumerate(svc_labels):
    sx = svc_start_x + i * (svc_w + svc_gap)
    box = add_rounded_rect(slide, sx, svc_top, svc_w, svc_h, WHITE)
    add_rect(slide, sx + Inches(0.05), svc_top, svc_w - Inches(0.1), Inches(0.04), color)
    ltb = add_text_box(slide, sx, svc_top + Inches(0.15), svc_w, svc_h - Inches(0.15))
    ltb.text_frame.word_wrap = True
    set_text(ltb.text_frame, label, size=10, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# API label
api_tb = add_text_box(slide, fe_left, gw_top + Inches(2.0), fe_w, Inches(0.4))
set_text(api_tb.text_frame, "15 REST API Controllers  |  JWT Auth  |  Role-Based Access", size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows down
arr_positions = [Inches(3.6), Inches(6.2), Inches(8.8)]
for ax in arr_positions:
    atb = add_text_box(slide, ax, gw_top + gw_h - Inches(0.1), Inches(1), Inches(0.5))
    set_text(atb.text_frame, "\u25BC", size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ── Bottom row: Three databases ──
db_top = Inches(5.8)
db_h = Inches(1.1)
db_w = Inches(2.3)
db_gap = Inches(0.35)
db_start = fe_left + Inches(0.2)

dbs = [
    ("PostgreSQL", "Geometry, Descriptors,\nFabrication, Jobs", ACCENT_GREEN),
    ("Apache Jena Fuseki", "Papers, Experiments,\nOutcomes (RDF/SPARQL)", ACCENT_AMBER),
    ("lovamap_core", "Computational Engine\n(Python)", ACCENT_TEAL),
]
for i, (name, desc, color) in enumerate(dbs):
    dx = db_start + i * (db_w + db_gap)
    db_box = add_rounded_rect(slide, dx, db_top, db_w, db_h, WHITE)
    add_rect(slide, dx + Inches(0.05), db_top, db_w - Inches(0.1), Inches(0.04), color)
    ntb = add_text_box(slide, dx, db_top + Inches(0.12), db_w, Inches(0.35))
    set_text(ntb.text_frame, name, size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    dtb = add_text_box(slide, dx, db_top + Inches(0.45), db_w, Inches(0.6))
    dtb.text_frame.word_wrap = True
    set_text(dtb.text_frame, desc, size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — PILLAR 1 DEEP DIVE: SCAFFOLD ANALYSIS
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_section_header(slide, 1, "Scaffold Analysis", ACCENT_TEAL)
add_slide_number(slide, 5)

# Left column: Job pipeline
left_tb = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5))
set_text(left_tb.text_frame, "Computational Job Pipeline", size=20, color=DUKE_NAVY, bold=True)

# Pipeline flow boxes
pipeline = [
    ("Upload", "CSV / DAT / JSON / TIF", ACCENT_TEAL),
    ("Segment", "Particle segmentation\nfrom microscopy", ACCENT_TEAL),
    ("Analyze", "Pore descriptors,\nglobal metrics", ACCENT_TEAL),
    ("Mesh", "3D .glb generation\nfor visualization", ACCENT_TEAL),
    ("Store", "Results persisted to\nScaffoldGroup + Descriptors", ACCENT_GREEN),
]

pipe_top = Inches(2.0)
pipe_w = Inches(5.2)
step_h = Inches(0.85)
step_gap = Inches(0.12)
for i, (name, desc, color) in enumerate(pipeline):
    sy = pipe_top + i * (step_h + step_gap)
    step_box = add_rounded_rect(slide, Inches(0.8), sy, pipe_w, step_h, WHITE)
    add_rect(slide, Inches(0.8), sy, Inches(0.08), step_h, color)

    num_tb = add_text_box(slide, Inches(1.1), sy + Inches(0.1), Inches(0.8), Inches(0.6))
    set_text(num_tb.text_frame, f"{i+1}.", size=20, color=color, bold=True)

    name_tb = add_text_box(slide, Inches(1.65), sy + Inches(0.05), Inches(1.2), Inches(0.4))
    set_text(name_tb.text_frame, name, size=15, color=DUKE_NAVY, bold=True)

    desc_tb = add_text_box(slide, Inches(1.65), sy + Inches(0.38), Inches(3.5), Inches(0.45))
    desc_tb.text_frame.word_wrap = True
    set_text(desc_tb.text_frame, desc, size=12, color=MID_GRAY)

    # Arrow between steps
    if i < len(pipeline) - 1:
        atb = add_text_box(slide, Inches(3.1), sy + step_h - Inches(0.05), Inches(0.5), Inches(0.3))
        set_text(atb.text_frame, "\u25BC", size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Right column: What we compute
right_tb = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.5))
set_text(right_tb.text_frame, "Computed Outputs", size=20, color=DUKE_NAVY, bold=True)

outputs = [
    ("Global Descriptors", "Scaffold-level scalar metrics: total porosity, mean pore volume, pore count, void fraction, surface area", ACCENT_TEAL),
    ("Pore Descriptors", "Per-pore measurement arrays: individual volumes, surface areas, aspect ratios, longest dimensions, connectivity", ACCENT_TEAL),
    ("Other Descriptors", "Additional computed properties: spatial distributions, throat metrics, coordination numbers", ACCENT_TEAL),
    ("3D Domains", "Segmented regions (particles, pores, other) as interactive .glb meshes with voxel-grid metadata", ACCENT_GREEN),
    ("AI Search", "Natural language queries across the scaffold database: 'spherical particles ~100\u00b5m with high porosity'", DUKE_BLUE),
]

for i, (title, desc, color) in enumerate(outputs):
    oy = Inches(2.0) + i * (Inches(0.95) + Inches(0.12))
    obox = add_rounded_rect(slide, Inches(6.8), oy, Inches(5.7), Inches(0.95), WHITE)
    add_rect(slide, Inches(6.8), oy, Inches(0.08), Inches(0.95), color)
    otb = add_text_box(slide, Inches(7.1), oy + Inches(0.08), Inches(5.2), Inches(0.35))
    set_text(otb.text_frame, title, size=14, color=DUKE_NAVY, bold=True)
    odtb = add_text_box(slide, Inches(7.1), oy + Inches(0.4), Inches(5.2), Inches(0.5))
    odtb.text_frame.word_wrap = True
    set_text(odtb.text_frame, desc, size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — PILLAR 2 DEEP DIVE: DATA STORAGE
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_section_header(slide, 2, "Comprehensive Data Storage", ACCENT_GREEN)
add_slide_number(slide, 6)

# Left half: Relational DB
left_title = add_text_box(slide, Inches(0.6), Inches(1.3), Inches(6), Inches(0.5))
set_text(left_title.text_frame, "Relational Database (PostgreSQL)", size=18, color=ACCENT_GREEN, bold=True)

left_sub = add_text_box(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.4))
set_text(left_sub.text_frame, "Scaffold instance data with full provenance", size=13, color=MID_GRAY)

# Data hierarchy visualization
hierarchy_text = """ScaffoldGroup
  \u251c\u2500 InputGroup
  \u2502    \u2514\u2500 ParticlePropertyGroups
  \u2502         (shape, size, stiffness, friction)
  \u251c\u2500 Scaffolds (replicates)
  \u2502    \u251c\u2500 Domains (.glb meshes)
  \u2502    \u251c\u2500 GlobalDescriptors (scalars)
  \u2502    \u251c\u2500 PoreDescriptors (arrays)
  \u2502    \u251c\u2500 OtherDescriptors (flexible)
  \u2502    \u2514\u2500 Jobs (analysis history)
  \u2514\u2500 Images & Tags"""

hier_box = add_rounded_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(3.7), WHITE)
hier_tb = add_text_box(slide, Inches(0.85), Inches(2.5), Inches(5.4), Inches(3.4))
hier_tb.text_frame.word_wrap = True
set_text(hier_tb.text_frame, hierarchy_text, size=13, color=DARK_TEXT, font_name="Courier New")

# Fabrication params note
fab_tb = add_text_box(slide, Inches(0.6), Inches(6.15), Inches(5.8), Inches(0.6))
fab_tb.text_frame.word_wrap = True
set_text(fab_tb.text_frame, "Supports batch upload (200 MB JSON), batch mesh replacement, and bulk descriptor seeding.", size=11, color=MID_GRAY)

# Divider
add_rect(slide, Inches(6.6), Inches(1.3), Inches(0.03), Inches(5.5), RGBColor(0xDD, 0xDD, 0xDD))

# Right half: RDF DB
right_title = add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.5))
set_text(right_title.text_frame, "RDF Knowledge Base (Fuseki)", size=18, color=ACCENT_AMBER, bold=True)

right_sub = add_text_box(slide, Inches(6.9), Inches(1.8), Inches(6), Inches(0.4))
set_text(right_sub.text_frame, "Paper-centric experimental knowledge", size=13, color=MID_GRAY)

rdf_classes = [
    ("Paper", "DOI, title, authors, journal, paper type"),
    ("Material", "Particle shape/size/material, porosity, packing"),
    ("FabricationMethod", "Chemistry, reagents (CAS#), enzymes, peptides"),
    ("Experiment", "Cell type, seeding density, culture conditions"),
    ("Outcome", "Value \u00b1 SD, unit, sample size, time series"),
    ("GeometryProfile", "Fuzzy match to relational DB scaffolds"),
]

for i, (cls, desc) in enumerate(rdf_classes):
    ry = Inches(2.3) + i * Inches(0.75)
    rbox = add_rounded_rect(slide, Inches(6.9), ry, Inches(5.8), Inches(0.65), WHITE)
    add_rect(slide, Inches(6.9), ry, Inches(0.06), Inches(0.65), ACCENT_AMBER)
    ctb = add_text_box(slide, Inches(7.2), ry + Inches(0.03), Inches(2.5), Inches(0.3))
    set_text(ctb.text_frame, cls, size=14, color=DUKE_NAVY, bold=True)
    ddtb = add_text_box(slide, Inches(7.2), ry + Inches(0.3), Inches(5.3), Inches(0.35))
    ddtb.text_frame.word_wrap = True
    set_text(ddtb.text_frame, desc, size=12, color=MID_GRAY)

# Bridge note
bridge_tb = add_text_box(slide, Inches(6.9), Inches(6.9), Inches(5.8), Inches(0.4))
bridge_tb.text_frame.word_wrap = True


# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — PILLAR 2 CONTINUED: BRIDGING THE DATABASES
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_section_header(slide, 2, "Bridging Relational & RDF Data", ACCENT_GREEN)
add_slide_number(slide, 7)

# Explanation
expl_tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8))
expl_tb.text_frame.word_wrap = True
set_text(expl_tb.text_frame,
    "The two databases serve complementary roles. The relational DB stores precise, per-scaffold computed data. "
    "The RDF store captures qualitative, paper-level knowledge. They connect through two bridging mechanisms:",
    size=15, color=DARK_TEXT)

# Left card: Exact match
left_card = add_rounded_rect(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.5), WHITE)
add_rect(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.06), ACCENT_GREEN)
ltb = add_text_box(slide, Inches(1.1), Inches(2.85), Inches(5.0), Inches(0.4))
set_text(ltb.text_frame, "Exact Match", size=20, color=ACCENT_GREEN, bold=True)
ltb2 = add_text_box(slide, Inches(1.1), Inches(3.4), Inches(5.0), Inches(2.5))
ltb2.text_frame.word_wrap = True
set_text(ltb2.text_frame, "lova:scaffoldGroupId  /  lova:scaffoldId", size=14, color=DUKE_NAVY, bold=True, font_name="Courier New")
add_paragraph(ltb2.text_frame, "", size=6, color=MID_GRAY)
add_paragraph(ltb2.text_frame, "Direct integer references from RDF Material nodes to relational DB records. Used when a paper's scaffold data has been uploaded to LOVAMAP.", size=14, color=MID_GRAY, space_before=Pt(8))
add_paragraph(ltb2.text_frame, "", size=6, color=MID_GRAY)
add_paragraph(ltb2.text_frame, 'Example: "Paper by Liu et al. (2023) used scaffolds that correspond to ScaffoldGroup #42 in our database."', size=13, color=DARK_TEXT, space_before=Pt(8))

# Right card: Fuzzy match
right_card = add_rounded_rect(slide, Inches(6.8), Inches(2.6), Inches(5.5), Inches(3.5), WHITE)
add_rect(slide, Inches(6.8), Inches(2.6), Inches(5.5), Inches(0.06), ACCENT_AMBER)
rtb = add_text_box(slide, Inches(7.1), Inches(2.85), Inches(5.0), Inches(0.4))
set_text(rtb.text_frame, "Fuzzy Match (GeometryProfile)", size=20, color=ACCENT_AMBER, bold=True)
rtb2 = add_text_box(slide, Inches(7.1), Inches(3.4), Inches(5.0), Inches(2.5))
rtb2.text_frame.word_wrap = True
set_text(rtb2.text_frame, "lova:hasGeometryProfile", size=14, color=DUKE_NAVY, bold=True, font_name="Courier New")
add_paragraph(rtb2.text_frame, "", size=6, color=MID_GRAY)
add_paragraph(rtb2.text_frame, "A GeometryProfile node describes matching criteria (particle shape, size range, material) to find similar scaffolds in the relational DB, even when the exact data hasn't been uploaded.", size=14, color=MID_GRAY, space_before=Pt(8))
add_paragraph(rtb2.text_frame, "", size=6, color=MID_GRAY)
add_paragraph(rtb2.text_frame, 'Example: "This paper used ~70\u00b5m PEG spheres \u2014 find all scaffold groups in our DB matching that profile."', size=13, color=DARK_TEXT, space_before=Pt(8))

# Bottom: what this enables
enables_tb = add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.7))
enables_tb.text_frame.word_wrap = True
set_text(enables_tb.text_frame, "This enables cross-referencing: ", size=15, color=DARK_TEXT, bold=True)
r = enables_tb.text_frame.paragraphs[0].add_run()
r.text = '"Show me the biological outcomes reported in the literature for scaffolds with geometry similar to the ones I just analyzed."'
r.font.size = Pt(15)
r.font.color.rgb = ACCENT_GREEN
r.font.italic = True
r.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — PILLAR 3: DATA STANDARDIZATION
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_section_header(slide, 3, "Data Standardization", ACCENT_AMBER)
add_slide_number(slide, 8)

# Four standardization aspects
aspects = [
    ("Controlled Descriptor Vocabulary",
     "Every computed property has a DescriptorType record defining its name, unit, category, data type, and description. When two scaffolds report 'pore volume,' they mean the same thing, computed the same way.",
     ACCENT_AMBER),
    ("Structured Fabrication Metadata",
     "InputGroup and ParticlePropertyGroup models enforce a consistent schema: container geometry, packing configuration (from controlled enums), per-particle-type properties with standardized fields. No more free-text fabrication descriptions.",
     ACCENT_AMBER),
    ("RDF Ontology as Reporting Template",
     "The LOVAMAP ontology defines the minimum information for a MAP scaffold study: materials (with CAS numbers), fabrication chemistry, experimental conditions (cell type, seeding density, culture duration), and outcomes (value \u00b1 SD, unit, sample size). Papers that fit this schema have reported with sufficient detail for reproducibility.",
     ACCENT_AMBER),
    ("Publication Dataset Versioning",
     "Frozen, hash-verified snapshots of scaffold data tied to specific papers. Descriptor selection rules control exactly which job outputs are included. The exact data underlying a publication can be reconstructed even as the live database evolves.",
     ACCENT_AMBER),
]

card_w = Inches(5.7)
card_h = Inches(2.4)
for i, (title, desc, color) in enumerate(aspects):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + col * (card_w + Inches(0.5))
    cy = Inches(1.4) + row * (card_h + Inches(0.3))

    card = add_rounded_rect(slide, cx, cy, card_w, card_h, WHITE)
    add_rect(slide, cx + Inches(0.1), cy, card_w - Inches(0.2), Inches(0.05), color)

    ttb = add_text_box(slide, cx + Inches(0.3), cy + Inches(0.2), card_w - Inches(0.6), Inches(0.4))
    set_text(ttb.text_frame, title, size=17, color=DUKE_NAVY, bold=True)

    dtb = add_text_box(slide, cx + Inches(0.3), cy + Inches(0.7), card_w - Inches(0.6), card_h - Inches(0.9))
    dtb.text_frame.word_wrap = True
    set_text(dtb.text_frame, desc, size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — VISION: WHERE WE'RE HEADED
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, DUKE_NAVY)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)
add_slide_number(slide, 9)

tb = add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Vision: Where We're Headed", size=32, color=WHITE, bold=True)

sub = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
set_text(sub.text_frame, "Building toward a comprehensive, community-driven resource for the MAP scaffold field.", size=15, color=RGBColor(0xB0, 0xC4, 0xDE))

vision_items = [
    ("\u25B6  Community database", " \u2014 A centralized, publicly accessible repository where any lab can upload scaffold data and browse results from across the field."),
    ("\u25B6  Literature knowledge graph", " \u2014 A growing RDF knowledge base that captures every MAP scaffold paper's materials, methods, and outcomes in a structured, queryable format."),
    ("\u25B6  Cross-study comparison", " \u2014 Overlay computed geometry from our relational DB with experimental outcomes from the literature to answer questions no single paper can."),
    ("\u25B6  Reporting standard", " \u2014 The LOVAMAP ontology as a community-endorsed minimum reporting checklist for MAP scaffold publications."),
    ("\u25B6  Automated ingestion", " \u2014 AI-assisted extraction of structured data from published papers, reducing the manual effort to populate the knowledge graph."),
    ("\u25B6  Expanded analysis", " \u2014 New descriptor types, mechanical simulation integration, and cross-scaffold comparison tools."),
]

vtb = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5.0))
vtb.text_frame.word_wrap = True
first = True
for bold_part, rest in vision_items:
    if first:
        p = vtb.text_frame.paragraphs[0]
        first = False
    else:
        p = vtb.text_frame.add_paragraph()
    p.space_before = Pt(14)
    p.space_after = Pt(4)
    r1 = p.add_run()
    r1.text = bold_part
    r1.font.size = Pt(17)
    r1.font.color.rgb = ACCENT_TEAL
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xD0, 0xD8, 0xE8)
    r2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — DISCUSSION: WHAT FUNCTIONALITY IS MOST USEFUL?
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CORAL)
add_slide_number(slide, 10)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Discussion: Functionality Priorities", size=30, color=DUKE_NAVY, bold=True)

sub = add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.5))
set_text(sub.text_frame, "What capabilities would make LOVAMAP most useful to the MAP scaffold community?", size=16, color=ACCENT_CORAL, bold=True)

# Two columns of discussion questions
col_w = Inches(5.7)
col1_left = Inches(0.7)
col2_left = Inches(6.9)

# Column 1: Analysis & tools
c1_title = add_text_box(slide, col1_left, Inches(1.8), col_w, Inches(0.5))
set_text(c1_title.text_frame, "Analysis & Tools", size=18, color=ACCENT_TEAL, bold=True)

q1_box = add_rounded_rect(slide, col1_left, Inches(2.3), col_w, Inches(4.7), WHITE)
q1_tb = add_text_box(slide, col1_left + Inches(0.3), Inches(2.5), col_w - Inches(0.6), Inches(4.3))
q1_tb.text_frame.word_wrap = True
questions1 = [
    "Which new descriptor types would be most valuable to compute? (e.g., mechanical properties, connectivity metrics, tortuosity)",
    "Should users be able to run custom analysis scripts on uploaded data, or is a curated set of descriptors sufficient?",
    "How important is cross-scaffold comparison tooling \u2014 e.g., overlaying pore distributions from multiple groups?",
    "Would simulation capabilities (e.g., mechanical loading, fluid flow) add meaningful value?",
    "What visualization features matter most: animations, cross-sections, overlays with experimental images?",
]
first = True
for q in questions1:
    if first:
        p = q1_tb.text_frame.paragraphs[0]
        first = False
    else:
        p = q1_tb.text_frame.add_paragraph()
    p.space_before = Pt(10)
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = f"\u25B8  {q}"
    r.font.size = Pt(13)
    r.font.color.rgb = DARK_TEXT
    r.font.name = "Calibri"

# Column 2: Data & community
c2_title = add_text_box(slide, col2_left, Inches(1.8), col_w, Inches(0.5))
set_text(c2_title.text_frame, "Data & Community", size=18, color=ACCENT_GREEN, bold=True)

q2_box = add_rounded_rect(slide, col2_left, Inches(2.3), col_w, Inches(4.7), WHITE)
q2_tb = add_text_box(slide, col2_left + Inches(0.3), Inches(2.5), col_w - Inches(0.6), Inches(4.3))
q2_tb.text_frame.word_wrap = True
questions2 = [
    "Should data upload be open to any researcher, or curated/moderated?",
    "How do we incentivize labs to contribute data? (credit, co-authorship, priority access?)",
    "What export formats are essential? (CSV, JSON, RDF/Turtle, custom reports?)",
    "How should we handle proprietary or pre-publication data? Embargo periods? Private groups?",
    "What level of metadata granularity is practical to require for uploads?",
]
first = True
for q in questions2:
    if first:
        p = q2_tb.text_frame.paragraphs[0]
        first = False
    else:
        p = q2_tb.text_frame.add_paragraph()
    p.space_before = Pt(10)
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = f"\u25B8  {q}"
    r.font.size = Pt(13)
    r.font.color.rgb = DARK_TEXT
    r.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
#  SLIDE 11 — DISCUSSION: WHAT QUERIES DO YOU EXPECT?
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CORAL)
add_slide_number(slide, 11)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Discussion: Expected Queries", size=30, color=DUKE_NAVY, bold=True)

sub = add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.6))
sub.text_frame.word_wrap = True
set_text(sub.text_frame, "When you imagine exploring this database, what questions would you want to answer?", size=16, color=ACCENT_CORAL, bold=True)

# Example queries we can already support
can_box = add_rounded_rect(slide, Inches(0.7), Inches(1.85), Inches(5.8), Inches(5.0), WHITE)
add_rect(slide, Inches(0.7), Inches(1.85), Inches(5.8), Inches(0.05), ACCENT_GREEN)
can_title = add_text_box(slide, Inches(1.0), Inches(2.05), Inches(5.2), Inches(0.4))
set_text(can_title.text_frame, "Queries We Can Support Today", size=17, color=ACCENT_GREEN, bold=True)

can_queries = [
    '"Show all scaffold groups with spherical particles between 50\u2013100\u00b5m"',
    '"What is the pore size distribution for ScaffoldGroup #42?"',
    '"Find scaffolds with porosity > 0.6 and isotropic packing"',
    '"Which papers studied PEG-based MAP scaffolds?"',
    '"What cell types have been tested on 70\u00b5m particle scaffolds?"',
    '"Show the knowledge graph for paper DOI:10.1002/adhm.202300823"',
    '"Compare pore aspect ratios across all my uploaded scaffold groups"',
]
ctb = add_text_box(slide, Inches(1.0), Inches(2.55), Inches(5.2), Inches(4.0))
ctb.text_frame.word_wrap = True
first = True
for q in can_queries:
    if first:
        p = ctb.text_frame.paragraphs[0]
        first = False
    else:
        p = ctb.text_frame.add_paragraph()
    p.space_before = Pt(7)
    r = p.add_run()
    r.text = f"\u2713  {q}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_TEXT
    r.font.name = "Calibri"

# Aspirational queries
asp_box = add_rounded_rect(slide, Inches(6.8), Inches(1.85), Inches(5.8), Inches(5.0), WHITE)
add_rect(slide, Inches(6.8), Inches(1.85), Inches(5.8), Inches(0.05), ACCENT_AMBER)
asp_title = add_text_box(slide, Inches(7.1), Inches(2.05), Inches(5.2), Inches(0.4))
set_text(asp_title.text_frame, "Queries We Want to Enable", size=17, color=ACCENT_AMBER, bold=True)

asp_queries = [
    '"What scaffold geometry maximizes cell infiltration for MSCs?"',
    '"Is there a correlation between pore aspect ratio and vascularization across all published studies?"',
    '"What fabrication parameters produce scaffolds most similar to this 3D-printed reference?"',
    '"Show me the optimal particle size range for bone regeneration based on all available data"',
    '"Which scaffold properties are most predictive of macrophage polarization?"',
    '"Compare my new scaffold\'s geometry to all published scaffolds with in vivo outcomes"',
]
atb = add_text_box(slide, Inches(7.1), Inches(2.55), Inches(5.2), Inches(4.0))
atb.text_frame.word_wrap = True
first = True
for q in asp_queries:
    if first:
        p = atb.text_frame.paragraphs[0]
        first = False
    else:
        p = atb.text_frame.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = f"\u25CB  {q}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_TEXT
    r.font.name = "Calibri"

# Bottom prompt
prompt_tb = add_text_box(slide, Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.4))
set_text(prompt_tb.text_frame, "What queries are missing? What questions would you ask this database that we haven't considered?",
    size=15, color=ACCENT_CORAL, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 12 — DISCUSSION: WHAT PROPERTIES TO PRIORITIZE?
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, SOFT_WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CORAL)
add_slide_number(slide, 12)

tb = add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Discussion: Schema Priorities", size=30, color=DUKE_NAVY, bold=True)

sub = add_text_box(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6))
sub.text_frame.word_wrap = True
set_text(sub.text_frame, "Which properties in our database schemas are most valuable to prioritize for standardizing the literature?", size=15, color=ACCENT_CORAL, bold=True)

# Three columns: Relational, RDF, New
col_w = Inches(3.8)
col_gap = Inches(0.35)
col_start = Inches(0.5)

# Column 1: Relational DB properties
c1_box = add_rounded_rect(slide, col_start, Inches(1.8), col_w, Inches(5.2), WHITE)
add_rect(slide, col_start, Inches(1.8), col_w, Inches(0.05), ACCENT_GREEN)
c1t = add_text_box(slide, col_start + Inches(0.2), Inches(1.95), col_w - Inches(0.4), Inches(0.4))
set_text(c1t.text_frame, "Relational DB (Geometry)", size=16, color=ACCENT_GREEN, bold=True)
c1sub = add_text_box(slide, col_start + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(0.3))
set_text(c1sub.text_frame, "Currently tracked:", size=12, color=MID_GRAY)

rel_props = [
    ("Particle:", " shape, size stats, stiffness, friction, dispersity, proportion"),
    ("Container:", " shape, size, dimensions"),
    ("Packing:", " isotropic, anisotropic, square, hexagonal"),
    ("Pore descriptors:", " volume, surface area, aspect ratio, longest length"),
    ("Global descriptors:", " porosity, pore count, mean metrics"),
    ("Domains:", " voxel size, mesh (.glb), segmentation source"),
]
c1tb = add_text_box(slide, col_start + Inches(0.2), Inches(2.65), col_w - Inches(0.4), Inches(3.0))
c1tb.text_frame.word_wrap = True
first = True
for bold, rest in rel_props:
    if first:
        p = c1tb.text_frame.paragraphs[0]
        first = False
    else:
        p = c1tb.text_frame.add_paragraph()
    p.space_before = Pt(5)
    r1 = p.add_run()
    r1.text = f"\u2022 {bold}"
    r1.font.size = Pt(12)
    r1.font.color.rgb = DUKE_NAVY
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(12)
    r2.font.color.rgb = MID_GRAY
    r2.font.name = "Calibri"

c1q = add_text_box(slide, col_start + Inches(0.2), Inches(5.3), col_w - Inches(0.4), Inches(1.5))
c1q.text_frame.word_wrap = True
set_text(c1q.text_frame, "What's missing?", size=13, color=ACCENT_CORAL, bold=True)
add_paragraph(c1q.text_frame, "\u25B8  Throat/window descriptors?", size=12, color=DARK_TEXT, space_before=Pt(6))
add_paragraph(c1q.text_frame, "\u25B8  Tortuosity / path length?", size=12, color=DARK_TEXT, space_before=Pt(4))
add_paragraph(c1q.text_frame, "\u25B8  Coordination number?", size=12, color=DARK_TEXT, space_before=Pt(4))
add_paragraph(c1q.text_frame, "\u25B8  Mechanical descriptors?", size=12, color=DARK_TEXT, space_before=Pt(4))

# Column 2: RDF properties
c2_left = col_start + col_w + col_gap
c2_box = add_rounded_rect(slide, c2_left, Inches(1.8), col_w, Inches(5.2), WHITE)
add_rect(slide, c2_left, Inches(1.8), col_w, Inches(0.05), ACCENT_AMBER)
c2t = add_text_box(slide, c2_left + Inches(0.2), Inches(1.95), col_w - Inches(0.4), Inches(0.4))
set_text(c2t.text_frame, "RDF (Literature & Experiments)", size=16, color=ACCENT_AMBER, bold=True)
c2sub = add_text_box(slide, c2_left + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(0.3))
set_text(c2sub.text_frame, "Currently tracked:", size=12, color=MID_GRAY)

rdf_props = [
    ("Material:", " shape, size, material, porosity, surface mods"),
    ("Fabrication:", " chemistry, reagents + CAS#, enzymes, peptides"),
    ("Experiment:", " cell type/source, seeding density, culture duration"),
    ("Environment:", " temperature, pH, atmosphere, medium"),
    ("Outcome:", " value \u00b1 SD, unit, sample size, time series"),
    ("Paper:", " DOI, authors, journal, focus material"),
]
c2tb = add_text_box(slide, c2_left + Inches(0.2), Inches(2.65), col_w - Inches(0.4), Inches(3.0))
c2tb.text_frame.word_wrap = True
first = True
for bold, rest in rdf_props:
    if first:
        p = c2tb.text_frame.paragraphs[0]
        first = False
    else:
        p = c2tb.text_frame.add_paragraph()
    p.space_before = Pt(5)
    r1 = p.add_run()
    r1.text = f"\u2022 {bold}"
    r1.font.size = Pt(12)
    r1.font.color.rgb = DUKE_NAVY
    r1.font.bold = True
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(12)
    r2.font.color.rgb = MID_GRAY
    r2.font.name = "Calibri"

c2q = add_text_box(slide, c2_left + Inches(0.2), Inches(5.3), col_w - Inches(0.4), Inches(1.5))
c2q.text_frame.word_wrap = True
set_text(c2q.text_frame, "What's missing?", size=13, color=ACCENT_CORAL, bold=True)
add_paragraph(c2q.text_frame, "\u25B8  Degradation profiles?", size=12, color=DARK_TEXT, space_before=Pt(6))
add_paragraph(c2q.text_frame, "\u25B8  In vivo outcome categories?", size=12, color=DARK_TEXT, space_before=Pt(4))
add_paragraph(c2q.text_frame, "\u25B8  Sterilization methods?",size=12, color=DARK_TEXT, space_before=Pt(4))
add_paragraph(c2q.text_frame, "\u25B8  Growth factor loading?", size=12, color=DARK_TEXT, space_before=Pt(4))

# Column 3: Reporting standard
c3_left = c2_left + col_w + col_gap
c3_box = add_rounded_rect(slide, c3_left, Inches(1.8), col_w, Inches(5.2), WHITE)
add_rect(slide, c3_left, Inches(1.8), col_w, Inches(0.05), ACCENT_CORAL)
c3t = add_text_box(slide, c3_left + Inches(0.2), Inches(1.95), col_w - Inches(0.4), Inches(0.4))
set_text(c3t.text_frame, "Reporting Standard", size=16, color=ACCENT_CORAL, bold=True)
c3sub = add_text_box(slide, c3_left + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(0.3))
set_text(c3sub.text_frame, "For the advisory board:", size=12, color=MID_GRAY)

c3tb = add_text_box(slide, c3_left + Inches(0.2), Inches(2.7), col_w - Inches(0.4), Inches(4.0))
c3tb.text_frame.word_wrap = True
report_qs = [
    "What should be the minimum required metadata for a scaffold upload?",
    "Which experimental details are essential vs. nice-to-have for cross-study comparison?",
    "Should we align with existing standards (e.g., MIAME, Dublin Core) or build our own?",
    "How granular should outcome reporting be? (mean only vs. raw data vs. individual timepoints?)",
    "What vocabulary/terminology should we enforce? (e.g., standard cell type names, standard assay names)",
    "How do we handle the gap between what papers report and what our schema requires?",
]

first = True
for q in report_qs:
    if first:
        p = c3tb.text_frame.paragraphs[0]
        first = False
    else:
        p = c3tb.text_frame.add_paragraph()
    p.space_before = Pt(7)
    p.space_after = Pt(2)
    r = p.add_run()
    r.text = f"\u25B8  {q}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_TEXT
    r.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
#  SLIDE 13 — KEY QUESTIONS SUMMARY
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, DUKE_NAVY)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)
add_slide_number(slide, 13)

tb = add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
set_text(tb.text_frame, "We'd Love Your Input On:", size=32, color=WHITE, bold=True)

key_questions = [
    ("1.", "Functionality", "What analysis tools or platform features would make LOVAMAP most useful to your work and the broader MAP scaffold community?"),
    ("2.", "Queries", "When exploring a database of MAP scaffold geometry and published experimental data, what are the first questions you'd want to answer?"),
    ("3.", "Relational Schema", "Which geometric, structural, or fabrication properties should we prioritize computing and storing for every scaffold?"),
    ("4.", "RDF Schema", "Which experimental details (cell types, assays, outcomes, fabrication chemistry) are most critical to capture from the literature?"),
    ("5.", "Reporting Standards", "What minimum information should be required when uploading scaffold data or ingesting a paper, to enable meaningful cross-study comparison?"),
]

for i, (num, title, desc) in enumerate(key_questions):
    qy = Inches(1.6) + i * Inches(1.1)

    # Number
    num_tb = add_text_box(slide, Inches(0.8), qy, Inches(0.5), Inches(0.5))
    set_text(num_tb.text_frame, num, size=22, color=ACCENT_TEAL, bold=True)

    # Title
    t_tb = add_text_box(slide, Inches(1.4), qy - Inches(0.03), Inches(2.5), Inches(0.4))
    set_text(t_tb.text_frame, title, size=18, color=WHITE, bold=True)

    # Description
    d_tb = add_text_box(slide, Inches(1.4), qy + Inches(0.35), Inches(11), Inches(0.6))
    d_tb.text_frame.word_wrap = True
    set_text(d_tb.text_frame, desc, size=14, color=RGBColor(0xB0, 0xC4, 0xDE))


# ═══════════════════════════════════════════════════════════════
#  SLIDE 14 — THANK YOU / CLOSING
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
fill_background(slide, DUKE_NAVY)
add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.08), ACCENT_TEAL)

tb = add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11), Inches(1.2))
set_text(tb.text_frame, "Thank You", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

sub = add_text_box(slide, Inches(1.0), Inches(3.5), Inches(11), Inches(0.8))
set_text(sub.text_frame, "We value your guidance in shaping LOVAMAP into a resource", size=20, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER)
add_paragraph(sub.text_frame, "the MAP scaffold community can build on together.", size=20, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER, space_before=Pt(2))

# Attribution
attr = add_text_box(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(1.0))
set_text(attr.text_frame, "Segura Lab  |  Duke University", size=16, color=RGBColor(0x8A, 0x9B, 0xAE), alignment=PP_ALIGN.CENTER)
add_paragraph(attr.text_frame, "lovamap.com", size=14, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER, space_before=Pt(8))


# ── Save ───────────────────────────────────────────────────────
output_path = "/Users/dzen/Documents/LOVAMAP/Repos/LOVAMAP_GW/LOVAMAP_Advisory_Board.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
